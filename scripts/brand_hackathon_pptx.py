#!/usr/bin/env python3
"""Aplikuje vizuální styl poslusnehlasim.cz na hackathonovou prezentaci."""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "Výstup - švejk.pptx"
SVEJK_SVG = ROOT / "svejk" / "static" / "svejk.svg"
SVEJK_PNG = ROOT / "svejk" / "static" / "svejk-terra.png"

# noviny-dlouhe.css
DESK = RGBColor(0xC9, 0xBD, 0xA3)
PAPER = RGBColor(0xF0, 0xE6, 0xCF)
INK = RGBColor(0x21, 0x1C, 0x14)
INK_SOFT = RGBColor(0x5A, 0x50, 0x3E)
TERRA = RGBColor(0xCF, 0x5A, 0x31)

FONT_DISPLAY = "Anton"
FONT_LABEL = "Oswald"
FONT_BODY = "Old Standard TT"


def _has_chrome(slide) -> bool:
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE.RECTANGLE:
            continue
        fill = sh.fill
        if fill.type == 1 and fill.fore_color.rgb == PAPER:  # MSO_FILL.SOLID
            return True
    return False


def _strip_chrome(slide) -> None:
    remove = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE.RECTANGLE:
            remove.append(sh)
            continue
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t in ("poslusnehlasim.cz", "ŠVEJK DO SNĚMOVNY"):
                remove.append(sh)
        if sh.shape_type == 13:  # PICTURE
            remove.append(sh)
    for sh in remove:
        sp = sh._element
        sp.getparent().remove(sp)


def _send_to_back(slide, shape) -> None:
    el = shape._element
    parent = el.getparent()
    parent.remove(el)
    slide.shapes._spTree.insert(2, el)


def _style_run(run, *, name: str, size: int, color: RGBColor, bold: bool = False) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def _style_paragraph_text(shape, body_size: int = 16) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        if not para.text.strip():
            continue
        for run in para.runs:
            _style_run(run, name=FONT_BODY, size=body_size, color=INK)


def _style_title_shape(shape) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        if not para.text.strip():
            continue
        para.text = para.text.upper()
        for run in para.runs:
            _style_run(run, name=FONT_DISPLAY, size=26, color=INK, bold=True)


def _ensure_svejk_png() -> Path:
    if SVEJK_PNG.is_file():
        return SVEJK_PNG
    if not SVEJK_SVG.is_file():
        return SVEJK_PNG
    SVEJK_PNG.parent.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "rsvg-convert",
            "-w",
            "200",
            "-o",
            str(SVEJK_PNG),
            str(SVEJK_SVG),
        ],
        check=True,
    )
    # terra silueta
    subprocess.run(
        [
            "magick",
            str(SVEJK_PNG),
            "-fuzz",
            "40%",
            "-fill",
            "#cf5a31",
            "-opaque",
            "black",
            "-alpha",
            "set",
            str(SVEJK_PNG),
        ],
        check=True,
    )
    return SVEJK_PNG


def _add_slide_chrome(slide, prs) -> None:
    if _has_chrome(slide):
        return
    w, h = prs.slide_width, prs.slide_height
    margin = Inches(0.28)
    bar_h = Inches(0.11)
    foot_h = Inches(0.28)

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DESK

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TERRA
    bar.line.fill.background()

    paper_top = bar_h + Inches(0.06)
    paper = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        margin,
        paper_top,
        w - 2 * margin,
        h - paper_top - foot_h,
    )
    paper.fill.solid()
    paper.fill.fore_color.rgb = PAPER
    paper.line.color.rgb = INK
    paper.line.width = Pt(1.25)

    rule_y = paper_top + Inches(0.04)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        margin,
        rule_y,
        w - 2 * margin,
        Pt(2),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = INK
    rule.line.fill.background()

    footer = slide.shapes.add_textbox(
        margin,
        h - foot_h + Inches(0.02),
        w - 2 * margin,
        foot_h,
    )
    ft = footer.text_frame
    ft.clear()
    p = ft.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "poslusnehlasim.cz"
    _style_run(r, name=FONT_LABEL, size=10, color=TERRA, bold=True)

    for sh in (bar, paper, rule):
        _send_to_back(slide, sh)


def _style_title_slide(slide) -> None:
    shapes = [s for s in slide.shapes if s.has_text_frame]
    if len(shapes) < 2:
        return
    title_sh, sub_sh = shapes[0], shapes[1]

    tf = title_sh.text_frame
    lines = [p.text for p in tf.paragraphs if p.text.strip()]
    subtitle_line = lines[1] if len(lines) > 1 else "Švejk do sněmovny"
    pitch = sub_sh.text_frame.text.strip()

    tf.clear()
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = "POSLUŠNĚ HLÁSÍM"
    _style_run(r0, name=FONT_DISPLAY, size=44, color=INK, bold=True)

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = subtitle_line
    _style_run(r1, name=FONT_LABEL, size=13, color=INK_SOFT)

    stf = sub_sh.text_frame
    stf.clear()
    p = stf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = pitch or (
        "Každý den z jednání PS dostanete noviny ve stylu Švejka: "
        "co se hlasovalo, proč na tom záleží a jak dlouho to v sále trvalo."
    )
    _style_run(r, name=FONT_BODY, size=18, color=INK)

    eyebrow = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(8.5), Inches(0.35))
    et = eyebrow.text_frame
    et.clear()
    ep = et.paragraphs[0]
    ep.alignment = PP_ALIGN.CENTER
    er = ep.add_run()
    er.text = "ŠVEJK DO SNĚMOVNY"
    _style_run(er, name=FONT_LABEL, size=11, color=TERRA, bold=True)

    png = _ensure_svejk_png()
    if png.is_file():
        slide.shapes.add_picture(
            str(png),
            Inches(7.15),
            Inches(0.75),
            height=Inches(1.55),
        )


def apply_brand(path: Path) -> None:
    prs = Presentation(str(path))
    for slide in prs.slides:
        _strip_chrome(slide)

    for idx, slide in enumerate(prs.slides):
        _add_slide_chrome(slide, prs)
        if idx == 0:
            _style_title_slide(slide)
            continue

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.is_placeholder:
                txt = shape.text_frame.text.strip()
                if txt == "poslusnehlasim.cz":
                    continue
            ph = shape.placeholder_format.type if shape.is_placeholder else None
            from pptx.enum.shapes import PP_PLACEHOLDER

            if ph == PP_PLACEHOLDER.TITLE or ph == PP_PLACEHOLDER.CENTER_TITLE:
                _style_title_shape(shape)
            elif ph == PP_PLACEHOLDER.BODY or ph == PP_PLACEHOLDER.SUBTITLE:
                _style_paragraph_text(shape, body_size=15 if idx == 6 else 14)
            else:
                _style_paragraph_text(shape, body_size=14)

    prs.save(str(path))
    print(f"Branded: {path}")


if __name__ == "__main__":
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else PPTX
    apply_brand(target)
